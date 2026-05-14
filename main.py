# PDF to Image Converter
# Copyright (c) 2026 Nate Rudquist
# All rights reserved.
# This software is proprietary and confidential.
# Unauthorized copying, distribution, or modification is prohibited.

import os
import sys
import shutil
import subprocess
import time
import random
import json
import urllib.request
import urllib.error
import urllib.parse
from contextlib import contextmanager
from html.parser import HTMLParser
from PySide6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QLabel, QPushButton, QProgressBar, QListWidget, QFileDialog,
    QLineEdit, QTabWidget, QTextEdit, QSpinBox, QMessageBox,
    QFrame, QScrollArea, QDialog, QDialogButtonBox,
)
from PySide6.QtCore import Qt, QThread, Signal
from PySide6.QtGui import QDragEnterEvent, QDropEvent
from pdf2image import convert_from_path, pdfinfo_from_path

# -- Locate Poppler ----------------------------------------------------------

def _find_poppler():
    if getattr(sys, "frozen", False):
        bundled = os.path.join(sys._MEIPASS, "poppler", "bin")
        if os.path.exists(bundled):
            return bundled
    if shutil.which("pdftoppm"):
        return str(os.path.dirname(shutil.which("pdftoppm")))
    return None

POPPLER_PATH = _find_poppler()

# Suppress black console windows spawned by Poppler subprocesses on Windows.
@contextmanager
def _no_window():
    if sys.platform != "win32":
        yield
        return
    orig = subprocess.Popen.__init__
    def _patched(self, *args, **kwargs):
        kwargs.setdefault("creationflags", 0)
        kwargs["creationflags"] |= subprocess.CREATE_NO_WINDOW
        orig(self, *args, **kwargs)
    subprocess.Popen.__init__ = _patched
    try:
        yield
    finally:
        subprocess.Popen.__init__ = orig

# -- Color palette -----------------------------------------------------------

BG         = "#1A1A1A"
SURFACE    = "#212121"
SURFACE2   = "#2A2A2A"
SURFACE3   = "#313131"
BORDER     = "#2E2E2E"
BORDER_LT  = "#3A3A3A"
GREEN      = "#4A7C59"
GREEN_LT   = "#5C9E6E"
GREEN_DIM  = "#2E5239"
TEXT       = "#E8E8E8"
TEXT_MUTED = "#6E6E6E"
TEXT_DIM   = "#9A9A9A"

BTN_H = 36   # standard control height

GLOBAL_STYLE = f"""
    QMainWindow, QWidget {{
        background-color: {BG};
        color: {TEXT};
        font-family: "Segoe UI";
        font-size: 13px;
    }}
    QTabWidget::pane {{
        border: 1px solid {BORDER};
        border-radius: 8px;
        background-color: {BG};
        top: -1px;
    }}
    QTabBar::tab {{
        background-color: {SURFACE};
        color: {TEXT_MUTED};
        border: 1px solid {BORDER};
        border-bottom: none;
        border-radius: 6px 6px 0 0;
        padding: 9px 24px;
        margin-right: 3px;
        font-size: 13px;
    }}
    QTabBar::tab:selected {{
        background-color: {BG};
        color: {TEXT};
        border-bottom: 1px solid {BG};
    }}
    QTabBar::tab:hover:!selected {{
        background-color: {SURFACE2};
        color: {TEXT_DIM};
    }}
    QListWidget {{
        background-color: {SURFACE};
        border: 1px solid {BORDER};
        border-radius: 6px;
        color: {TEXT};
        padding: 4px;
        outline: none;
    }}
    QListWidget::item {{
        padding: 5px 8px;
        border-radius: 4px;
    }}
    QListWidget::item:selected {{
        background-color: {GREEN_DIM};
        color: {TEXT};
    }}
    QListWidget::item:hover {{
        background-color: {SURFACE2};
    }}
    QTextEdit {{
        background-color: {SURFACE};
        border: 1px solid {BORDER};
        border-radius: 6px;
        color: {TEXT};
        font-size: 12px;
        padding: 6px;
    }}
    QSpinBox {{
        background-color: {SURFACE};
        border: 1px solid {BORDER};
        border-radius: 6px;
        color: {TEXT};
        padding: 0 8px;
        font-size: 13px;
        height: {BTN_H}px;
    }}
    QSpinBox::up-button, QSpinBox::down-button {{ width: 0; }}
    QProgressBar {{
        background-color: {SURFACE};
        border: 1px solid {BORDER};
        border-radius: 5px;
        height: 6px;
        text-align: center;
        color: transparent;
    }}
    QProgressBar::chunk {{
        background-color: {GREEN};
        border-radius: 4px;
    }}
    QPushButton {{
        background-color: {SURFACE2};
        color: {TEXT_DIM};
        border: 1px solid {BORDER_LT};
        border-radius: 6px;
        padding: 0 18px;
        height: {BTN_H}px;
        font-size: 13px;
    }}
    QPushButton:hover {{
        background-color: {SURFACE3};
        color: {TEXT};
        border-color: #484848;
    }}
    QPushButton:pressed {{
        background-color: {SURFACE};
    }}
    QPushButton:disabled {{
        color: #3A3A3A;
        border-color: #252525;
        background-color: {SURFACE};
    }}
    QScrollArea {{
        border: none;
        background-color: transparent;
    }}
    QScrollBar:vertical {{
        background: {SURFACE};
        width: 8px;
        border-radius: 4px;
        margin: 0;
    }}
    QScrollBar::handle:vertical {{
        background: {BORDER_LT};
        border-radius: 4px;
        min-height: 24px;
    }}
    QScrollBar::handle:vertical:hover {{
        background: #4A4A4A;
    }}
    QScrollBar::add-line:vertical, QScrollBar::sub-line:vertical {{ height: 0; }}
"""

ACTION_BTN_STYLE = f"""
    QPushButton {{
        background-color: {GREEN};
        color: #FFFFFF;
        border: none;
        border-radius: 6px;
        padding: 0 28px;
        height: {BTN_H}px;
        font-size: 13px;
        font-weight: 600;
    }}
    QPushButton:hover {{
        background-color: {GREEN_LT};
    }}
    QPushButton:pressed {{
        background-color: {GREEN_DIM};
    }}
    QPushButton:disabled {{
        background-color: #253529;
        color: #3D5C44;
        border: none;
    }}
"""

FOLDER_INPUT_STYLE = f"""
    QLineEdit {{
        background-color: {SURFACE};
        border: 1px solid {BORDER};
        border-radius: 6px;
        color: {TEXT_MUTED};
        font-size: 12px;
        padding: 0 10px;
        height: {BTN_H}px;
    }}
"""

URL_INPUT_STYLE = f"""
    QLineEdit {{
        background-color: {SURFACE};
        border: 1px solid {BORDER};
        border-radius: 6px;
        color: {TEXT};
        font-size: 13px;
        padding: 0 10px;
        height: {BTN_H}px;
    }}
    QLineEdit:focus {{
        border-color: {GREEN_DIM};
    }}
"""

SECTION_HDR_STYLE  = f"color: {TEXT}; font-size: 12px; font-weight: 600;"
HINT_STYLE         = f"color: {TEXT_MUTED}; font-size: 11px;"
STATUS_STYLE       = f"color: {TEXT_DIM}; font-size: 12px;"

# ---------------------------------------------------------------------------
# Workers
# ---------------------------------------------------------------------------

class ConvertWorker(QThread):
    progress     = Signal(int)
    status       = Signal(str)
    error        = Signal(str)
    images_ready = Signal(list)
    finished     = Signal()

    def __init__(self, pdf_paths, output_folder):
        super().__init__()
        self.pdf_paths = pdf_paths
        self.output_folder = output_folder

    def run(self):
        total_pages = 0
        page_counts = []
        for path in self.pdf_paths:
            try:
                with _no_window():
                    info = pdfinfo_from_path(path, poppler_path=POPPLER_PATH)
                page_counts.append(info["Pages"])
                total_pages += info["Pages"]
            except Exception:
                filename = os.path.basename(path)
                self.error.emit(
                    f"Could not read '{filename}' — file may be corrupted, encrypted, or not a valid PDF."
                )
                page_counts.append(0)

        try:
            os.makedirs(self.output_folder, exist_ok=True)
        except Exception:
            self.error.emit(
                "Could not create the output folder — check that you have permission to write there."
            )
            self.finished.emit()
            return

        completed_pages = 0
        created_images = []
        skipped = 0
        for pdf_path, num_pages in zip(self.pdf_paths, page_counts):
            if num_pages == 0:
                skipped += 1
                continue
            filename = os.path.basename(pdf_path)
            stem = os.path.splitext(filename)[0]

            for page_num in range(1, num_pages + 1):
                self.status.emit(
                    f"Converting {filename}  —  page {page_num} of {num_pages}"
                )
                try:
                    with _no_window():
                        images = convert_from_path(
                            pdf_path,
                            first_page=page_num,
                            last_page=page_num,
                            poppler_path=POPPLER_PATH,
                        )
                    out_name = f"{stem}.jpg" if num_pages == 1 else f"{stem}_page_{page_num}.jpg"
                    out_path = os.path.join(self.output_folder, out_name)
                    images[0].save(out_path, "JPEG")
                    created_images.append(out_path)
                except PermissionError:
                    self.error.emit(
                        f"Could not write to the output folder — check that you have permission to write there."
                    )
                    self.images_ready.emit(created_images)
                    self.finished.emit()
                    return
                except Exception:
                    self.error.emit(
                        f"Could not convert '{filename}' page {page_num} — file may be corrupted or encrypted."
                    )
                completed_pages += 1
                if total_pages > 0:
                    self.progress.emit(int(completed_pages / total_pages * 100))

        self.images_ready.emit(created_images)
        if skipped and not created_images:
            self.status.emit("Conversion failed — see errors above.")
        elif skipped:
            self.status.emit(f"Done with {skipped} skipped file(s). Images saved to the output folder.")
        else:
            self.status.emit("Done!  Images saved to the output folder.")

        self.finished.emit()

class ScanWorker(QThread):
    found    = Signal(list)
    status   = Signal(str)
    finished = Signal()

    def __init__(self, url):
        super().__init__()
        self.url = url

    def run(self):
        class LinkParser(HTMLParser):
            def __init__(self, base):
                super().__init__()
                self.base = base
                self.links = []

            def handle_starttag(self, tag, attrs):
                if tag == "a":
                    for attr, val in attrs:
                        if attr == "href" and val and ".pdf" in val.lower():
                            full = urllib.parse.urljoin(self.base, val)
                            if full not in self.links:
                                self.links.append(full)

        self.status.emit(f"Scanning {self.url}…")
        try:
            req = urllib.request.Request(
                self.url, headers={"User-Agent": "Mozilla/5.0"}
            )
            with urllib.request.urlopen(req, timeout=15) as r:
                html = r.read().decode("utf-8", errors="ignore")
            parser = LinkParser(self.url)
            parser.feed(html)
            self.found.emit(parser.links)
            self.status.emit(f"Found {len(parser.links)} PDF link(s).")
        except urllib.error.HTTPError as e:
            if e.code == 403:
                self.status.emit(
                    "This page requires authentication. Download the files through "
                    "your browser instead, then drag them into the Convert tab."
                )
            else:
                self.status.emit(f"Scan error: HTTP {e.code}")
            self.found.emit([])
        except Exception as e:
            self.status.emit(f"Scan error: {e}")
            self.found.emit([])
        self.finished.emit()

class DownloadWorker(QThread):
    progress    = Signal(int)
    status      = Signal(str)
    log         = Signal(str)
    paths_ready = Signal(list)
    finished    = Signal()

    def __init__(self, urls, output_folder, log_file):
        super().__init__()
        self.urls = urls
        self.output_folder = output_folder
        self.log_file = log_file
        self._stop = False

    def stop(self):
        self._stop = True

    def _load_log(self):
        if os.path.exists(self.log_file):
            with open(self.log_file, "r") as f:
                return set(json.load(f))
        return set()

    def _save_log(self, downloaded):
        with open(self.log_file, "w") as f:
            json.dump(sorted(downloaded), f, indent=2)

    BASE_HEADERS = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/124.0.0.0 Safari/537.36"
        ),
        "Accept": "application/pdf,text/html,application/xhtml+xml;q=0.9,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.9",
        "Accept-Encoding": "gzip, deflate, br",
        "Connection": "keep-alive",
        "Upgrade-Insecure-Requests": "1",
    }

    def run(self):
        os.makedirs(self.output_folder, exist_ok=True)
        downloaded = self._load_log()
        saved_paths = []

        total = len(self.urls)
        completed = 0

        for url in self.urls:
            if self._stop:
                self.log.emit("— Stopped by user.")
                break

            filename = url.split("?")[0].rstrip("/").split("/")[-1]
            if not filename.lower().endswith(".pdf"):
                filename += ".pdf"

            filepath = os.path.join(self.output_folder, filename)

            if url in downloaded:
                self.log.emit(f"  skip  {filename}")
                if os.path.exists(filepath):
                    saved_paths.append(filepath)
                completed += 1
                self.progress.emit(int(completed / total * 100))
                continue

            self.status.emit(f"Downloading {filename}…")
            try:
                req = urllib.request.Request(url, headers=dict(self.BASE_HEADERS))
                with urllib.request.urlopen(req, timeout=30) as r:
                    data = r.read()

                with open(filepath, "wb") as f:
                    f.write(data)
                downloaded.add(url)
                self._save_log(downloaded)
                saved_paths.append(filepath)
                self.log.emit(f"  ✓  {filename}")

            except urllib.error.HTTPError as e:
                if e.code == 403:
                    self.log.emit(
                        f"  ✗  {filename}  — Access denied. Download this file manually "
                        "in your browser, then drag it into the Convert tab."
                    )
                elif e.code == 404:
                    self.log.emit(f"  ✗  {filename}  — File not found (404). The URL may be incorrect or the file removed.")
                else:
                    self.log.emit(f"  ✗  {filename}  — Download failed (HTTP {e.code}). Check your internet connection.")
            except urllib.error.URLError:
                self.log.emit(f"  ✗  {filename}  — Could not connect. Check your internet connection.")
            except PermissionError:
                self.log.emit(f"  ✗  {filename}  — Could not save file — check that you have permission to write to the download folder.")
            except Exception:
                self.log.emit(f"  ✗  {filename}  — Download failed. Check your internet connection and try again.")

            completed += 1
            self.progress.emit(int(completed / total * 100))

            if completed < total and not self._stop:
                time.sleep(random.uniform(1.0, 2.0))

        self.status.emit("Finished.")
        self.paths_ready.emit(saved_paths)
        self.finished.emit()

# ---------------------------------------------------------------------------
# PowerPoint export worker
# ---------------------------------------------------------------------------

class PptxWorker(QThread):
    status   = Signal(str)
    finished = Signal(int, list)   # (file_count, [saved_paths])

    SLIDES_PER_FILE = 200

    def __init__(self, image_paths, output_folder, base_name):
        super().__init__()
        self.image_paths = image_paths
        self.output_folder = output_folder
        self.base_name = base_name

    def run(self):
        try:
            from pptx import Presentation
            from pptx.util import Emu
        except ImportError:
            self.finished.emit(0, [])
            self.status.emit(
                "python-pptx is not installed. Run: pip install python-pptx"
            )
            return

        SLIDE_W = Emu(9_144_000)   # 10 inches
        SLIDE_H = Emu(6_858_000)   # 7.5 inches
        chunk_size = self.SLIDES_PER_FILE
        chunks = [
            self.image_paths[i : i + chunk_size]
            for i in range(0, len(self.image_paths), chunk_size)
        ]
        saved = []

        for idx, chunk in enumerate(chunks):
            prs = Presentation()
            prs.slide_width  = SLIDE_W
            prs.slide_height = SLIDE_H
            blank = prs.slide_layouts[6]

            for img_path in chunk:
                self.status.emit(f"Adding {os.path.basename(img_path)}…")
                slide = prs.slides.add_slide(blank)
                slide.shapes.add_picture(img_path, 0, 0, SLIDE_W, SLIDE_H)

            fname = (
                f"{self.base_name}.pptx"
                if len(chunks) == 1
                else f"{self.base_name}_Part{idx + 1}.pptx"
            )
            out_path = os.path.join(self.output_folder, fname)
            self.status.emit(f"Saving {fname}…")
            prs.save(out_path)
            saved.append(out_path)

        self.finished.emit(len(saved), saved)

# ---------------------------------------------------------------------------
# Drop zone widget
# ---------------------------------------------------------------------------

class DropZone(QLabel):
    files_dropped = Signal(list)

    def __init__(self):
        super().__init__("Drop PDFs here\nor click to browse")
        self.setAlignment(Qt.AlignCenter)
        self.setAcceptDrops(True)
        self.setMinimumHeight(130)
        self._set_style(active=False)

    def _set_style(self, active):
        if active:
            bg, border, color = GREEN_DIM, GREEN_LT, GREEN_LT
        else:
            bg, border, color = SURFACE, BORDER_LT, TEXT_MUTED
        self.setStyleSheet(f"""
            QLabel {{
                background-color: {bg};
                border: 2px dashed {border};
                border-radius: 8px;
                color: {color};
                font-size: 14px;
                font-family: "Segoe UI";
            }}
        """)

    def dragEnterEvent(self, event: QDragEnterEvent):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
            self._set_style(active=True)

    def dragLeaveEvent(self, event):
        self._set_style(active=False)

    def dropEvent(self, event: QDropEvent):
        paths = [
            url.toLocalFile()
            for url in event.mimeData().urls()
            if url.toLocalFile().lower().endswith(".pdf")
        ]
        self._set_style(active=False)
        if paths:
            self.files_dropped.emit(paths)

    def mousePressEvent(self, event):
        self.files_dropped.emit([])

# ---------------------------------------------------------------------------
# Export dialog
# ---------------------------------------------------------------------------

class ExportDialog(QDialog):
    def __init__(self, default_folder, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Export to PowerPoint")
        self.setMinimumWidth(460)
        self._folder = default_folder

        layout = QVBoxLayout(self)
        layout.setSpacing(12)
        layout.setContentsMargins(24, 24, 24, 20)

        # Filename
        name_lbl = QLabel("Filename")
        name_lbl.setStyleSheet(SECTION_HDR_STYLE)
        layout.addWidget(name_lbl)

        self.name_edit = QLineEdit("Slides")
        self.name_edit.setStyleSheet(URL_INPUT_STYLE)
        self.name_edit.setPlaceholderText("e.g. Slides")
        layout.addWidget(self.name_edit)

        hint = QLabel("Do not include the .pptx extension. For 200+ images, _Part1, _Part2 … will be appended automatically.")
        hint.setWordWrap(True)
        hint.setStyleSheet(HINT_STYLE)
        layout.addWidget(hint)

        layout.addSpacing(4)

        # Save location
        loc_lbl = QLabel("Save location")
        loc_lbl.setStyleSheet(SECTION_HDR_STYLE)
        layout.addWidget(loc_lbl)

        loc_row = QHBoxLayout()
        loc_row.setSpacing(8)
        self.folder_display = QLineEdit(default_folder)
        self.folder_display.setReadOnly(True)
        self.folder_display.setStyleSheet(FOLDER_INPUT_STYLE)
        browse_btn = QPushButton("Browse…")
        browse_btn.clicked.connect(self._pick_folder)
        loc_row.addWidget(self.folder_display)
        loc_row.addWidget(browse_btn)
        layout.addLayout(loc_row)

        layout.addSpacing(8)

        # OK / Cancel
        buttons = QDialogButtonBox(QDialogButtonBox.Ok | QDialogButtonBox.Cancel)
        buttons.button(QDialogButtonBox.Ok).setText("Export")
        buttons.accepted.connect(self._accept)
        buttons.rejected.connect(self.reject)
        layout.addWidget(buttons)

    def _pick_folder(self):
        folder = QFileDialog.getExistingDirectory(self, "Choose Save Location", self._folder)
        if folder:
            self._folder = folder
            self.folder_display.setText(folder)

    def _accept(self):
        name = self.name_edit.text().strip()
        # Strip .pptx if user typed it
        if name.lower().endswith(".pptx"):
            name = name[:-5].strip()
        if not name:
            self.name_edit.setFocus()
            return
        self.name_edit.setText(name)
        self.accept()

    def result_values(self):
        return self.name_edit.text().strip(), self._folder

# ---------------------------------------------------------------------------
# Main window
# ---------------------------------------------------------------------------

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("PDF to Image v1.1")
        self.setMinimumSize(660, 680)

        if getattr(sys, "frozen", False):
            self.base_dir = os.path.dirname(sys.executable)
        else:
            self.base_dir = os.path.dirname(os.path.abspath(__file__))

        self.convert_output  = os.path.join(self.base_dir, "output")
        self.download_output = os.path.join(self.base_dir, "output", "downloads")
        self.pdf_paths = []
        self.last_downloaded_paths = []
        self.session_images = []

        self._build_ui()

        if not POPPLER_PATH:
            self.convert_status.setText(
                "Warning: Poppler not found. Conversion will fail. "
                "Please install Poppler and add it to your PATH."
            )

    # -- Layout helpers ------------------------------------------------------

    def _divider(self):
        line = QFrame()
        line.setFixedHeight(1)
        line.setStyleSheet(f"background: {BORDER}; border: none;")
        return line

    def _show_about(self):
        dlg = QDialog(self)
        dlg.setWindowTitle("About PDF to Image Converter")
        dlg.setMinimumWidth(400)

        layout = QVBoxLayout(dlg)
        layout.setSpacing(10)
        layout.setContentsMargins(24, 24, 24, 20)

        name = QLabel("PDF to Image Converter")
        name.setStyleSheet(f"font-size: 17px; font-weight: 700; color: {TEXT};")
        layout.addWidget(name)

        version = QLabel("Version 1.1")
        version.setStyleSheet(f"color: {TEXT_DIM}; font-size: 12px;")
        layout.addWidget(version)

        layout.addSpacing(4)

        desc = QLabel(
            "A tool for downloading PDFs from web pages and converting them "
            "to images for research purposes."
        )
        desc.setWordWrap(True)
        desc.setStyleSheet(f"color: {TEXT}; font-size: 13px;")
        layout.addWidget(desc)

        layout.addSpacing(4)

        copy = QLabel("Copyright \u00a9 2026 Nate Rudquist. All rights reserved.")
        copy.setStyleSheet(f"color: {TEXT_DIM}; font-size: 12px;")
        copy.setWordWrap(True)
        layout.addWidget(copy)

        prop = QLabel("This software is proprietary and confidential.")
        prop.setStyleSheet(f"color: {TEXT_MUTED}; font-size: 11px;")
        layout.addWidget(prop)

        layout.addSpacing(6)

        disclaimer = QLabel(
            "This software is provided AS-IS with no warranty or ongoing support. "
            "Use at your own risk."
        )
        disclaimer.setWordWrap(True)
        disclaimer.setStyleSheet(f"color: {TEXT_MUTED}; font-size: 11px;")
        layout.addWidget(disclaimer)

        layout.addSpacing(8)

        buttons = QDialogButtonBox(QDialogButtonBox.Ok)
        buttons.accepted.connect(dlg.accept)
        layout.addWidget(buttons)

        dlg.exec()

    def _section_label(self, text):
        lbl = QLabel(text)
        lbl.setStyleSheet(SECTION_HDR_STYLE)
        return lbl

    def _build_ui(self):
        help_menu = self.menuBar().addMenu("Help")
        about_action = help_menu.addAction("About PDF to Image Converter")
        about_action.triggered.connect(self._show_about)

        central = QWidget()
        self.setCentralWidget(central)
        root = QVBoxLayout(central)
        root.setContentsMargins(20, 16, 20, 20)
        root.setSpacing(12)

        title = QLabel("PDF to Image")
        title.setStyleSheet(
            f"font-size: 22px; font-weight: 700; color: {TEXT};"
        )
        root.addWidget(title)

        self.tabs = QTabWidget()
        self.tabs.addTab(self._build_convert_tab(), "Convert")
        self.tabs.addTab(self._build_download_tab(), "Download")
        root.addWidget(self.tabs)

    # -- Convert tab ---------------------------------------------------------

    def _build_convert_tab(self):
        tab = QWidget()
        layout = QVBoxLayout(tab)
        layout.setSpacing(0)
        layout.setContentsMargins(20, 20, 20, 20)

        # — Drop zone + file list —
        self.drop_zone = DropZone()
        self.drop_zone.files_dropped.connect(self._on_files_dropped)
        layout.addWidget(self.drop_zone)

        layout.addSpacing(10)

        self.file_list = QListWidget()
        self.file_list.setFixedHeight(110)
        layout.addWidget(self.file_list)

        layout.addSpacing(12)

        # — Clear / Convert buttons —
        btn_row = QHBoxLayout()
        btn_row.setSpacing(10)
        self.clear_btn = QPushButton("Clear List")
        self.clear_btn.clicked.connect(self._clear)
        self.convert_btn = QPushButton("Convert to Images")
        self.convert_btn.setStyleSheet(ACTION_BTN_STYLE)
        self.convert_btn.clicked.connect(self._start_conversion)
        self.convert_btn.setEnabled(False)
        btn_row.addWidget(self.clear_btn)
        btn_row.addStretch()
        btn_row.addWidget(self.convert_btn)
        layout.addLayout(btn_row)

        layout.addSpacing(16)
        layout.addWidget(self._divider())
        layout.addSpacing(14)

        # — Output folder —
        layout.addWidget(self._section_label("Output folder"))
        layout.addSpacing(8)

        layout.addLayout(self._folder_row("convert"))

        layout.addSpacing(16)
        layout.addWidget(self._divider())
        layout.addSpacing(14)

        # — Progress & status —
        self.convert_progress = QProgressBar()
        self.convert_progress.setValue(0)
        layout.addWidget(self.convert_progress)

        layout.addSpacing(8)

        self.convert_status = QLabel("Drop PDFs above to get started.")
        self.convert_status.setAlignment(Qt.AlignCenter)
        self.convert_status.setStyleSheet(STATUS_STYLE)
        self.convert_status.setWordWrap(True)
        layout.addWidget(self.convert_status)

        layout.addSpacing(16)
        layout.addWidget(self._divider())
        layout.addSpacing(14)

        # — Export —
        layout.addWidget(self._section_label("Export"))
        layout.addSpacing(8)

        self.export_pptx_btn = QPushButton("Export to PowerPoint")
        self.export_pptx_btn.setStyleSheet(ACTION_BTN_STYLE)
        self.export_pptx_btn.clicked.connect(self._export_to_pptx)
        layout.addWidget(self.export_pptx_btn)

        layout.addStretch()
        return tab

    # -- Download tab --------------------------------------------------------

    def _build_download_tab(self):
        tab = QWidget()
        tab_layout = QVBoxLayout(tab)
        tab_layout.setContentsMargins(0, 0, 0, 0)
        tab_layout.setSpacing(0)

        # Wrap everything in a scroll area so sections breathe
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setHorizontalScrollBarPolicy(Qt.ScrollBarAlwaysOff)
        scroll.setFrameShape(QFrame.Shape.NoFrame)

        inner = QWidget()
        layout = QVBoxLayout(inner)
        layout.setSpacing(0)
        layout.setContentsMargins(20, 20, 20, 20)

        scroll.setWidget(inner)
        tab_layout.addWidget(scroll)

        # — Scan page —
        layout.addWidget(self._section_label("Scan page for PDF links"))
        layout.addSpacing(8)

        scan_row = QHBoxLayout()
        scan_row.setSpacing(8)
        self.scan_input = QLineEdit()
        self.scan_input.setPlaceholderText("Paste a webpage URL…")
        self.scan_input.setStyleSheet(URL_INPUT_STYLE)
        self.scan_btn = QPushButton("Scan Page")
        self.scan_btn.clicked.connect(self._scan_page)
        scan_row.addWidget(self.scan_input)
        scan_row.addWidget(self.scan_btn)
        layout.addLayout(scan_row)

        layout.addSpacing(6)

        self.scan_status = QLabel("")
        self.scan_status.setStyleSheet(HINT_STYLE)
        self.scan_status.setWordWrap(True)
        layout.addWidget(self.scan_status)

        layout.addSpacing(18)
        layout.addWidget(self._divider())
        layout.addSpacing(16)

        # — URL generator —
        layout.addWidget(self._section_label("Generate URL list"))
        layout.addSpacing(8)

        gen_row = QHBoxLayout()
        gen_row.setSpacing(8)

        self.gen_pattern = QLineEdit()
        self.gen_pattern.setPlaceholderText("URL pattern with {n}  e.g.  https://example.com/file{n}.pdf")
        self.gen_pattern.setStyleSheet(URL_INPUT_STYLE)

        start_lbl = QLabel("Start:")
        start_lbl.setStyleSheet(f"color: {TEXT_DIM}; font-size: 12px;")
        self.gen_start = QSpinBox()
        self.gen_start.setRange(1, 9_999_999)
        self.gen_start.setValue(1)
        self.gen_start.setFixedWidth(72)

        end_lbl = QLabel("End:")
        end_lbl.setStyleSheet(f"color: {TEXT_DIM}; font-size: 12px;")
        self.gen_end = QSpinBox()
        self.gen_end.setRange(1, 9_999_999)
        self.gen_end.setValue(100)
        self.gen_end.setFixedWidth(72)

        pad_lbl = QLabel("Pad:")
        pad_lbl.setStyleSheet(f"color: {TEXT_DIM}; font-size: 12px;")
        self.gen_pad = QSpinBox()
        self.gen_pad.setRange(0, 12)
        self.gen_pad.setValue(0)
        self.gen_pad.setFixedWidth(52)
        self.gen_pad.setToolTip("Zero-pad numbers to this many digits (0 = none)")

        gen_btn = QPushButton("Generate")
        gen_btn.clicked.connect(self._generate_urls)

        gen_row.addWidget(self.gen_pattern)
        gen_row.addWidget(start_lbl)
        gen_row.addWidget(self.gen_start)
        gen_row.addWidget(end_lbl)
        gen_row.addWidget(self.gen_end)
        gen_row.addWidget(pad_lbl)
        gen_row.addWidget(self.gen_pad)
        gen_row.addWidget(gen_btn)
        layout.addLayout(gen_row)

        layout.addSpacing(6)

        gen_hint = QLabel("Use {n} as the number placeholder — e.g. file{n}.pdf with Pad 8 → file00000001.pdf")
        gen_hint.setStyleSheet(HINT_STYLE)
        gen_hint.setWordWrap(True)
        layout.addWidget(gen_hint)

        layout.addSpacing(18)
        layout.addWidget(self._divider())
        layout.addSpacing(16)

        # — URL list —
        layout.addWidget(self._section_label("PDF URLs to download"))
        layout.addSpacing(8)

        self.url_box = QTextEdit()
        self.url_box.setPlaceholderText(
            "Paste PDF URLs here, one per line.\n"
            "Or use Scan Page above to fill this automatically."
        )
        self.url_box.setFixedHeight(120)
        layout.addWidget(self.url_box)

        layout.addSpacing(18)
        layout.addWidget(self._divider())
        layout.addSpacing(16)

        # — Save folder + download controls —
        layout.addWidget(self._section_label("Save location"))
        layout.addSpacing(8)

        self.clear_downloads_btn = QPushButton("Clear Downloads Folder")
        self.clear_downloads_btn.clicked.connect(self._clear_downloads_folder)
        layout.addLayout(self._folder_row("download", extra_btn=self.clear_downloads_btn))

        layout.addSpacing(12)

        dl_btn_row = QHBoxLayout()
        dl_btn_row.setSpacing(10)
        self.stop_btn = QPushButton("Stop")
        self.stop_btn.setEnabled(False)
        self.stop_btn.clicked.connect(self._stop_download)
        self.download_btn = QPushButton("Download")
        self.download_btn.setStyleSheet(ACTION_BTN_STYLE)
        self.download_btn.clicked.connect(self._start_download)
        dl_btn_row.addWidget(self.stop_btn)
        dl_btn_row.addStretch()
        dl_btn_row.addWidget(self.download_btn)
        layout.addLayout(dl_btn_row)

        layout.addSpacing(12)

        self.download_progress = QProgressBar()
        self.download_progress.setValue(0)
        layout.addWidget(self.download_progress)

        layout.addSpacing(8)

        self.download_status = QLabel("Paste URLs or scan a page to get started.")
        self.download_status.setAlignment(Qt.AlignCenter)
        self.download_status.setStyleSheet(STATUS_STYLE)
        self.download_status.setWordWrap(True)
        layout.addWidget(self.download_status)

        layout.addSpacing(18)
        layout.addWidget(self._divider())
        layout.addSpacing(16)

        # — Download log —
        layout.addWidget(self._section_label("Log"))
        layout.addSpacing(8)

        self.log_view = QTextEdit()
        self.log_view.setReadOnly(True)
        self.log_view.setPlaceholderText("Download activity will appear here…")
        self.log_view.setFixedHeight(110)
        self.log_view.setStyleSheet(
            f"QTextEdit {{ font-family: 'Consolas', monospace; font-size: 11px; "
            f"background: {SURFACE}; border: 1px solid {BORDER}; border-radius: 6px; }}"
        )
        layout.addWidget(self.log_view)

        layout.addSpacing(12)

        self.convert_dl_btn = QPushButton("Convert Downloaded Files")
        self.convert_dl_btn.setStyleSheet(ACTION_BTN_STYLE)
        self.convert_dl_btn.setEnabled(False)
        self.convert_dl_btn.clicked.connect(self._convert_downloaded)
        layout.addWidget(self.convert_dl_btn)

        layout.addStretch()
        return tab

    # -- Shared folder-picker row --------------------------------------------

    def _folder_row(self, which, extra_btn=None):
        row = QHBoxLayout()
        row.setSpacing(8)

        lbl = QLabel("Save to:")
        lbl.setStyleSheet(f"color: {TEXT_DIM}; font-size: 12px;")
        lbl.setFixedWidth(58)

        folder = self.convert_output if which == "convert" else self.download_output
        display = QLineEdit(folder)
        display.setReadOnly(True)
        display.setStyleSheet(FOLDER_INPUT_STYLE)

        browse = QPushButton("Browse…")

        if which == "convert":
            self.convert_folder_display = display
            browse.clicked.connect(lambda: self._pick_folder("convert"))
        else:
            self.download_folder_display = display
            browse.clicked.connect(lambda: self._pick_folder("download"))

        row.addWidget(lbl)
        row.addWidget(display)
        row.addWidget(browse)
        if extra_btn is not None:
            row.addWidget(extra_btn)
        return row

    def _pick_folder(self, which):
        current = self.convert_output if which == "convert" else self.download_output
        folder = QFileDialog.getExistingDirectory(self, "Choose Output Folder", current)
        if not folder:
            return
        if which == "convert":
            self.convert_output = folder
            self.convert_folder_display.setText(folder)
        else:
            self.download_output = folder
            self.download_folder_display.setText(folder)

    # -- Convert logic -------------------------------------------------------

    def _on_files_dropped(self, paths):
        if not paths:
            paths, _ = QFileDialog.getOpenFileNames(
                self, "Select PDFs", "", "PDF Files (*.pdf)"
            )
        for path in paths:
            if path not in self.pdf_paths:
                self.pdf_paths.append(path)
                self.file_list.addItem(os.path.basename(path))
        self.convert_btn.setEnabled(bool(self.pdf_paths))

    def _clear(self):
        self.pdf_paths.clear()
        self.file_list.clear()
        self.convert_progress.setValue(0)
        self.convert_status.setText("Drop PDFs above to get started.")
        self.convert_btn.setEnabled(False)

    def _start_conversion(self):
        self.convert_btn.setEnabled(False)
        self.clear_btn.setEnabled(False)
        self.convert_progress.setValue(0)
        self.convert_status.setText("Starting…")

        self.session_images = []
        self.convert_worker = ConvertWorker(self.pdf_paths, self.convert_output)
        self.convert_worker.progress.connect(self.convert_progress.setValue)
        self.convert_worker.status.connect(self.convert_status.setText)
        self.convert_worker.error.connect(
            lambda e: self.convert_status.setText(f"Error: {e}")
        )
        self.convert_worker.images_ready.connect(self._on_images_ready)
        self.convert_worker.finished.connect(self._on_convert_finished)
        self.convert_worker.start()

    def _on_images_ready(self, images):
        self.session_images = images

    def _on_convert_finished(self):
        self.convert_btn.setEnabled(True)
        self.clear_btn.setEnabled(True)

    def _clear_downloads_folder(self):
        folder = self.download_output
        if not os.path.exists(folder):
            self.download_status.setText("Downloads folder does not exist.")
            return
        files = [f for f in os.listdir(folder) if os.path.isfile(os.path.join(folder, f))]
        if not files:
            self.download_status.setText("Downloads folder is already empty.")
            return
        reply = QMessageBox.question(
            self,
            "Clear Downloads Folder",
            f"Delete all {len(files)} file(s) in the downloads folder? This cannot be undone.",
            QMessageBox.Yes | QMessageBox.No,
            QMessageBox.No,
        )
        if reply == QMessageBox.Yes:
            deleted = 0
            for fname in files:
                try:
                    os.remove(os.path.join(folder, fname))
                    deleted += 1
                except Exception:
                    pass
            self.download_status.setText(f"Deleted {deleted} file(s) from the downloads folder.")

    def _export_to_pptx(self):
        images = [p for p in self.session_images if os.path.exists(p)]
        if not images:
            self.convert_status.setText(
                "No images from the current session. Convert PDFs first."
            )
            return

        dlg = ExportDialog(self.convert_output, parent=self)
        if dlg.exec() != QDialog.Accepted:
            return
        base_name, folder = dlg.result_values()

        self.export_pptx_btn.setEnabled(False)
        self.convert_status.setText(f"Building PowerPoint from {len(images)} image(s)…")
        self.pptx_worker = PptxWorker(images, folder, base_name)
        self.pptx_worker.status.connect(self.convert_status.setText)
        self.pptx_worker.finished.connect(self._on_pptx_finished)
        self.pptx_worker.start()

    def _on_pptx_finished(self, file_count, paths):
        self.export_pptx_btn.setEnabled(True)
        if file_count == 0:
            return
        names = ", ".join(os.path.basename(p) for p in paths)
        self.convert_status.setText(f"Saved {file_count} PowerPoint file(s): {names}")

    # -- Scan logic ----------------------------------------------------------

    def _generate_urls(self):
        pattern = self.gen_pattern.text().strip()
        if not pattern or "{n}" not in pattern:
            self.scan_status.setText("Pattern must contain {n} — e.g. https://example.com/file{n}.pdf")
            return
        start = self.gen_start.value()
        end   = self.gen_end.value()
        pad   = self.gen_pad.value()
        if start > end:
            self.scan_status.setText("Start must be ≤ End.")
            return
        urls = [pattern.replace("{n}", str(i).zfill(pad) if pad else str(i))
                for i in range(start, end + 1)]
        existing = self.url_box.toPlainText().strip()
        combined = (existing + "\n" if existing else "") + "\n".join(urls)
        self.url_box.setPlainText(combined.strip())
        self.scan_status.setText(f"Added {len(urls)} URL(s) to the list.")

    def _scan_page(self):
        url = self.scan_input.text().strip()
        if not url:
            return
        self.scan_btn.setEnabled(False)
        self.scan_status.setText("Scanning…")

        self.scan_worker = ScanWorker(url)
        self.scan_worker.found.connect(self._on_scan_found)
        self.scan_worker.status.connect(self.scan_status.setText)
        self.scan_worker.finished.connect(lambda: self.scan_btn.setEnabled(True))
        self.scan_worker.start()

    def _on_scan_found(self, links):
        if not links:
            return
        existing = self.url_box.toPlainText().strip()
        existing_set = set(existing.splitlines()) if existing else set()
        new_links = [l for l in links if l not in existing_set]
        if new_links:
            combined = (existing + "\n" if existing else "") + "\n".join(new_links)
            self.url_box.setPlainText(combined.strip())

    # -- Download logic ------------------------------------------------------

    def _start_download(self):
        raw = self.url_box.toPlainText().strip()
        urls = [u.strip() for u in raw.splitlines() if u.strip().startswith("http")]
        if not urls:
            self.download_status.setText("No valid URLs found. Paste URLs or scan a page first.")
            return

        log_file = os.path.join(self.download_output, "download_log.json")

        self.download_btn.setEnabled(False)
        self.stop_btn.setEnabled(True)
        self.convert_dl_btn.setEnabled(False)
        self.download_progress.setValue(0)
        self.log_view.clear()
        self.log_view.append(f"Starting {len(urls)} download(s)…\n")

        self.dl_worker = DownloadWorker(urls, self.download_output, log_file)
        self.dl_worker.progress.connect(self.download_progress.setValue)
        self.dl_worker.status.connect(self.download_status.setText)
        self.dl_worker.log.connect(self._append_log)
        self.dl_worker.paths_ready.connect(self._on_paths_ready)
        self.dl_worker.finished.connect(self._on_download_finished)
        self.dl_worker.start()

    def _stop_download(self):
        if hasattr(self, "dl_worker"):
            self.dl_worker.stop()
        self.stop_btn.setEnabled(False)

    def _append_log(self, text):
        self.log_view.append(text)
        self.log_view.verticalScrollBar().setValue(
            self.log_view.verticalScrollBar().maximum()
        )

    def _on_paths_ready(self, paths):
        self.last_downloaded_paths = paths

    def _on_download_finished(self):
        self.download_btn.setEnabled(True)
        self.stop_btn.setEnabled(False)
        self.convert_dl_btn.setEnabled(bool(self.last_downloaded_paths))

    # -- Convert downloaded files --------------------------------------------

    def _convert_downloaded(self):
        paths = [p for p in self.last_downloaded_paths if os.path.exists(p)]
        if not paths:
            return

        self._clear()
        for path in paths:
            self.pdf_paths.append(path)
            self.file_list.addItem(os.path.basename(path))
        self.convert_btn.setEnabled(True)

        self.tabs.setCurrentIndex(0)
        self._start_conversion()

# ---------------------------------------------------------------------------

if __name__ == "__main__":
    app = QApplication(sys.argv)
    app.setStyleSheet(GLOBAL_STYLE)
    window = MainWindow()
    window.show()
    sys.exit(app.exec())

# ---------------------------------------------------------------------------
# PDF to Image Converter
# Copyright (c) 2026 Nate Rudquist
# All rights reserved.
# This software is proprietary and confidential.
# Unauthorized copying, distribution, or modification is prohibited.
# ---------------------------------------------------------------------------
